import os
import re
from docx import Document
import docx2txt
from pathlib import Path
import win32com.client
import pythoncom
import time
from tqdm import tqdm
import json
import PyPDF2
import fitz  # PyMuPDF
import chardet
import pptx  # python-pptx

class TextSegmenter:
    def __init__(self, text_dir="source_article", output_dir="segmented_output"):

        self.text_dir = Path(text_dir)
        self.output_dir = Path(output_dir)
        self.output_dir.mkdir(exist_ok=True)
        
    def _doc_to_docx(self, doc_path):
        docx_path = doc_path.with_suffix('.docx')
        
        if docx_path.exists():
            return docx_path
            
        print(f"转换 {doc_path} 为 DOCX 格式...")
        try:
            # 初始化COM对象
            pythoncom.CoInitialize()
            word = win32com.client.Dispatch("Word.Application")
            word.Visible = False
            
            # 打开文档
            doc = word.Documents.Open(str(doc_path.absolute()))
            # 保存为docx
            doc.SaveAs(str(docx_path.absolute()), 16)  # 16表示docx格式
            doc.Close()
            word.Quit()
            
            return docx_path
        except Exception as e:
            print(f"转换文档时出错: {e}")
            return None
    
    def _ppt_to_pptx(self, ppt_path):
        """将PPT文件转换为PPTX格式"""
        pptx_path = ppt_path.with_suffix('.pptx')
        
        if pptx_path.exists():
            return pptx_path
            
        print(f"转换 {ppt_path} 为 PPTX 格式...")
        try:
            # 初始化COM对象
            pythoncom.CoInitialize()
            powerpoint = win32com.client.Dispatch("PowerPoint.Application")
            powerpoint.Visible = True
            
            # 打开PPT
            presentation = powerpoint.Presentations.Open(str(ppt_path.absolute()))
            # 保存为pptx
            presentation.SaveAs(str(pptx_path.absolute()), 24)  # 24表示pptx格式
            presentation.Close()
            powerpoint.Quit()
            
            return pptx_path
        except Exception as e:
            print(f"转换PPT文档时出错: {e}")
            return None
    
    def _rtf_to_txt(self, rtf_path):
        """将RTF文件转换为纯文本"""
        txt_path = rtf_path.with_suffix('.txt')
        
        if txt_path.exists():
            return txt_path
            
        print(f"转换 {rtf_path} 为 TXT 格式...")
        try:
            # 初始化COM对象
            pythoncom.CoInitialize()
            word = win32com.client.Dispatch("Word.Application")
            word.Visible = False
            
            # 打开文档
            doc = word.Documents.Open(str(rtf_path.absolute()))
            # 保存为txt
            doc.SaveAs(str(txt_path.absolute()), 2)  # 2表示txt格式
            doc.Close()
            word.Quit()
            
            return txt_path
        except Exception as e:
            print(f"转换RTF文档时出错: {e}")
            return None
    
    def _extract_from_pdf(self, pdf_path):
        """从PDF文件中提取文本"""
        try:
            # 尝试使用PyMuPDF (fitz)提取文本，通常效果更好
            try:
                text = []
                doc = fitz.open(pdf_path)
                for page in doc:
                    text.append(page.get_text())
                return "\n".join(text)
            except Exception as e:
                print(f"使用PyMuPDF提取PDF失败: {e}，尝试使用PyPDF2...")
                
            # 备选方案：使用PyPDF2
            with open(pdf_path, 'rb') as file:
                reader = PyPDF2.PdfReader(file)
                text = []
                for page in reader.pages:
                    text.append(page.extract_text() or "")
                return "\n".join(text)
        except Exception as e:
            print(f"从PDF提取文本时出错: {e}")
            return ""
    
    def _extract_from_txt(self, txt_path):
        """从TXT文件中提取文本，自动检测编码"""
        try:
            # 首先检测文件编码
            with open(txt_path, 'rb') as f:
                raw_data = f.read()
                result = chardet.detect(raw_data)
                encoding = result['encoding'] or 'utf-8'
            
            # 使用检测到的编码打开文件
            with open(txt_path, 'r', encoding=encoding, errors='replace') as f:
                return f.read()
        except Exception as e:
            print(f"从TXT文件提取文本时出错: {e}")
            return ""
    
    def _extract_from_pptx(self, pptx_path):
        """从PPTX文件中提取文本"""
        try:
            presentation = pptx.Presentation(pptx_path)
            text_content = []
            
            # 提取幻灯片标题和内容
            for i, slide in enumerate(presentation.slides):
                slide_text = []
                
                # 提取标题
                if slide.shapes.title and slide.shapes.title.text:
                    slide_text.append(f"## 幻灯片 {i+1}: {slide.shapes.title.text}")
                else:
                    slide_text.append(f"## 幻灯片 {i+1}")
                
                # 提取所有文本框内容
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text and shape != slide.shapes.title:
                        slide_text.append(shape.text)
                
                text_content.append("\n".join(slide_text))
            
            return "\n\n".join(text_content)
        except Exception as e:
            print(f"从PPTX提取文本时出错: {e}")
            return ""
    
    def extract_text_from_file(self, file_path):
        """从各种文件格式中提取文本"""
        file_path = Path(file_path)
        suffix = file_path.suffix.lower()
        title = file_path.stem
        
        # 提取文本内容
        content = ""
        
        # 根据文件类型选择不同的处理方法
        if suffix == '.doc':
            docx_path = self._doc_to_docx(file_path)
            if docx_path is None:
                print(f"无法从 {file_path} 提取文本")
                return ""
            file_path = docx_path
            suffix = '.docx'
        
        elif suffix == '.ppt':
            pptx_path = self._ppt_to_pptx(file_path)
            if pptx_path is None:
                print(f"无法从 {file_path} 提取文本")
                return ""
            file_path = pptx_path
            suffix = '.pptx'
        
        if suffix == '.docx':
            try:
                doc = Document(file_path)
                full_text = []
                
                # 提取正文段落
                for para in doc.paragraphs:
                    if para.text.strip():
                        full_text.append(para.text)
                
                content = "\n".join(full_text)
            except Exception as e:
                print(f"读取DOCX文件 {file_path} 时出错: {e}")
                # 尝试使用docx2txt作为备选方案
                try:
                    content = docx2txt.process(file_path)
                except:
                    content = ""
        
        elif suffix == '.pptx':
            content = self._extract_from_pptx(file_path)
        
        elif suffix == '.pdf':
            content = self._extract_from_pdf(file_path)
            
        elif suffix == '.txt':
            content = self._extract_from_txt(file_path)
            
        elif suffix == '.rtf':
            txt_path = self._rtf_to_txt(file_path)
            if txt_path:
                content = self._extract_from_txt(txt_path)
        
        else:
            print(f"不支持的文件格式: {suffix}")
            return ""
        
        # 添加标题作为第一行
        if content:
            content = f"# {title}\n\n" + content
            
        return content
    
    def segment_by_paragraph(self, text, min_length=50):
        """按段落分块，过滤掉太短的段落"""
        paragraphs = text.split('\n')
        return [p for p in paragraphs if len(p.strip()) >= min_length]
            
    def segment_by_heading(self, text):
        """按标题（章节、小节）分块"""
        # 使用正则表达式匹配常见的标题格式
        # 例如：第一章、1.1、## 标题、Title等
        heading_patterns = [
            r'第[一二三四五六七八九十零〇百千万亿\d]+[章节篇].*?(?=第[一二三四五六七八九十零〇百千万亿\d]+[章节篇]|$)',
            r'\d+\.\d+.*?(?=\d+\.\d+|$)',
            r'#+\s.*?(?=#+\s|$)',
            r'[一二三四五六七八九十]+、.*?(?=[一二三四五六七八九十]+、|$)'
        ]
        
        chunks = []
        remaining_text = text
        
        # 尝试所有标题模式，累积匹配结果
        for pattern in heading_patterns:
            pattern_matches = list(re.finditer(pattern, remaining_text, re.DOTALL))
            
            if pattern_matches:
                for match in pattern_matches:
                    chunks.append(match.group(0).strip())
        
        # 如果没有找到标题结构，退回到段落分块
        if not chunks:
            chunks = self.segment_by_paragraph(text)
        
        return chunks
    
    def process_files(self):
        
        result = {}
        
        # 扩展支持的文件类型
        supported_extensions = ['.doc', '.docx', '.pdf', '.txt', '.rtf', '.ppt', '.pptx']
        doc_files = []
        
        for ext in supported_extensions:
            doc_files.extend(list(self.text_dir.glob(f'*{ext}')))
        
        for file_path in tqdm(doc_files, desc="处理文件"):
            file_name = file_path.stem
            print(f"\n处理文件: {file_name}")
            
            # 提取文本
            text = self.extract_text_from_file(file_path)
            if not text:
                continue
                
            # 分块
            chunks = self.segment_by_heading(text)
            
            # 保存结果
            result[file_name] = {
                'file_path': str(file_path),
                'file_type': file_path.suffix,
                'chunk_count': len(chunks),
                'chunks': chunks
            }
            
            # 将每个文件的分块结果保存为单独的JSON文件
            output_file = self.output_dir / f"{file_name}.json"
            with open(output_file, 'w', encoding='utf-8') as f:
                json.dump({
                    'file_name': file_name,
                    'file_type': file_path.suffix,
                    'segmentation_method': 'heading',
                    'chunk_count': len(chunks),
                    'chunks': chunks
                }, f, ensure_ascii=False, indent=2)
                
            print(f"已分块: {len(chunks)} 个片段，保存至 {output_file}")
        
            
        return result

